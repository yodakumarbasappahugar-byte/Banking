import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def build_final_ppt():
    prs = Presentation()
    
    # Image Paths (Current Conversation)
    base_dir = r"C:\Users\HP\.gemini\antigravity\brain\c4094568-2eac-4a0d-bc02-bba6e1a83a86"
    hero_img = os.path.join(base_dir, "banking_hero_glassmorphism_1775462563278.png")
    security_img = os.path.join(base_dir, "security_lock_neon_1775462626055.png")
    tech_img = os.path.join(base_dir, "tech_stack_nodes_1775462642742.png")

    def set_bg(slide, color=RGBColor(10, 14, 22)): # Deep Midnight
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def apply_title_style(shape, color=RGBColor(99, 102, 241)): # Indigo Accent
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.size = Pt(36)
                paragraph.font.color.rgb = color
                paragraph.alignment = PP_ALIGN.LEFT

    def add_slide(layout_idx, title_text):
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        set_bg(slide)
        if slide.shapes.title:
            slide.shapes.title.text = title_text
            apply_title_style(slide.shapes.title)
        return slide

    def add_bullets(slide, points):
        if len(slide.placeholders) < 2: return
        tf = slide.placeholders[1].text_frame
        tf.word_wrap = True
        for point in points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(226, 232, 240)
            p.space_after = Pt(10)

    # 1. Title Slide
    s1 = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    set_bg(s1)
    if os.path.exists(hero_img):
        s1.shapes.add_picture(hero_img, Inches(0), Inches(0), width=Inches(5.5), height=Inches(7.5))
    
    # Title Box
    title_box = s1.shapes.add_textbox(Inches(5.8), Inches(2.5), Inches(4), Inches(2))
    tf = title_box.text_frame
    p1 = tf.add_paragraph()
    p1.text = "NIDHI BANK"
    p1.font.bold = True
    p1.font.size = Pt(54)
    p1.font.color.rgb = RGBColor(99, 102, 241)
    
    p2 = tf.add_paragraph()
    p2.text = "Premium Digital Banking Ecosystem"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(255, 255, 255)

    data = [
        ("Executive Project Summary", [
            "Vision: A high-fidelity, secure, and modern fintech solution.",
            "Design: Implementation of the 'Glassmorphism' aesthetic (translucency & blur).",
            "Functional: Real-time fund transfers, virtual card system, and administrative hub.",
            "Architecture: Fully decoupled Frontend-Backend-Database orchestration.",
            "Goal: To provide a zero-friction banking experience for modern users."
        ]),
        ("Modern Tech Stack Architecture", [
            "Frontend: Next.js 14+ (React) utilizing App Router and CSS Modules.",
            "Backend: FastAPI (Python) - High performance, asynchronous data handling.",
            "Database: Neon PostgreSQL - Serverless relational storage for ACID integrity.",
            "DevOps: Automated deployment via Vercel (Frontend) and Render (Backend).",
            "Communication: Secure RESTful API interactions with global CORS hardening."
        ]),
        ("Premium UI/UX: Glassmorphism", [
            "Aesthetic: Blurred translucent panels with vibrant glowing accents.",
            "Color Palette: Midnight blue foundations with Indigo and Violet highlights.",
            "Interactivity: Smooth frame-based animations and pulsing notification badges.",
            "Responsiveness: 100% mobile-adaptive navigation with toggleable sidebars.",
            "Localization: Standardized Indian Rupee (₹) symbol and en-IN formatting."
        ]),
        ("Backend Engineering: FastAPI", [
            "Security: Bcrypt salting and hashing for credential protection.",
            "Validation: Pydantic-driven data serialization and error handling.",
            "Database Connectivity: psycopg2 with RealDictCursor for efficient mapping.",
            "Performance: Asynchronous endpoint handling for high-concurrency requests.",
            "Integrity: Server-side validation of all financial transfer operations."
        ]),
        ("Database: Neon PostgreSQL", [
            "Data Tables: Normalized schema for 'Users' and 'Transactions'.",
            "Financial Precision: Using DECIMAL(12,2) for exact monetary calculations.",
            "Relationships: Foreign key constraints linking ledgers to user identities.",
            "Auditability: Global timestamps ensuring a robust transaction audit trail.",
            "Architecture: Scalable serverless branching for development environments."
        ]),
        ("Financial Logic: Fund Transfers", [
            "Atomicity: Single-transaction SQL blocks for simultaneous balance updates.",
            "Integrity: Automatic 'Rollback' on any failure during transfer processing.",
            "Validation: Server-side balance checks to prevent overdrafts.",
            "History: Immediate logging of transfers into the immutable ledger.",
            "Interface: Quick-transfer form with real-time success feedback."
        ]),
        ("Security Hub & Defense", [
            "Credentials: Zero-knowledge password storage using production Bcrypt.",
            "Network: Rigid CORS policy management on the FastAPI Render node.",
            "Visibility: Security center with login activity and device monitoring.",
            "Prevention: Mitigation of legacy hashing vulnerabilities (max-length limits).",
            "Access: Multi-layer validation for profile identity access."
        ]),
        ("Administrative User Management", [
            "Management: Centralized directory for all registered banking customers.",
            "Search: Real-time filtering by name, email, or account number.",
            "Monitoring: Visibility into active balances, Member IDs, and IFSC codes.",
            "Onboarding: Ability to manage 10+ mock accounts for system demonstration.",
            "Dashboard: Integrated overview of recent system-wide activity."
        ]),
        ("Card Management Ecosystem", [
            "Virtual Cards: Interactive digital debit/credit card generation.",
            "Controls: Real-time tools to 'Freeze Card' or 'Reset PIN'.",
            "Security: Mock security credentials allowing for safe testing.",
            "Visuals: Glassmorphism card UI with glowing border indicators.",
            "Integration: Full linkage between card transactions and user balance."
        ]),
        ("Recent Innovation Milestones", [
            "Real-time Sync: Fixed stale balance issues in the Account Identity modal.",
            "Event Signaling: Linked UI refreshes to the 'new-transaction' event.",
            "Redeploy Automation: Implemented 'redeploy.py' for 1-click CI/CD updates.",
            "Localization: Completed global transition to Indian Rupee (₹) standard.",
            "Stability: Enhanced error handling for cross-origin database requests."
        ]),
        ("Deployment & Final Vision", [
            "Frontend Hosting: Vercel (Production URL: nidhibank.vercel.app).",
            "Backend Web Service: Render (Service ID: banking-backend-api).",
            "CI/CD: Automated Git synchronization across environments.",
            "Scalability: Foundation ready for real-world user scaling.",
            "Summary: A complete, functional Fintech application delivered with excellence."
        ]),
        ("Questions & Discussion", [
            "Thank you for your attention!",
            "Project Repository: NidhiBank (Private GitHub)",
            "Architecture: Next.js + FastAPI + PostgreSQL",
            "Developed with Antigravity AI Support",
            "Any Questions?"
        ])
    ]

    for title, points in data:
        slide = add_slide(1, title)
        add_bullets(slide, points)
        
        # Add supporting images to specific slides
        if "Tech Stack" in title and os.path.exists(tech_img):
            slide.shapes.add_picture(tech_img, Inches(7.5), Inches(4.5), width=Inches(2))
        if "Security" in title and os.path.exists(security_img):
            slide.shapes.add_picture(security_img, Inches(7.5), Inches(2.2), width=Inches(2))

    prs.save("NidhiBank_Final_Detailed_Presentation.pptx")
    print("Presentation created successfully: NidhiBank_Final_Detailed_Presentation.pptx")

if __name__ == "__main__":
    build_final_ppt()
